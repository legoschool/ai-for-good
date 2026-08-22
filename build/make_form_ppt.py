# -*- coding: utf-8 -*-
"""선생님이 채워 넣는 양식 PPTX 를 만든다.

두 가지를 담는다.
  1. 오류 제보  : 화면 캡처를 붙이고 무슨 일이 있었는지 적는다
  2. 수업 후기  : 수업 사진을 붙이고 어땠는지 적는다

캔바로 가져가 나눠 쓰는 것을 전제로 만든다.
그래서 글상자는 겹치지 않게 두고, 사진 자리는 굵은 테두리 상자로 비워 둔다.

사용법 : py -3 build/make_form_ppt.py
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

import tasks as T

T.setup_console()

FONT = "맑은 고딕"
W, H = Inches(13.333), Inches(7.5)

INK = RGBColor(0x11, 0x11, 0x11)
MUTED = RGBColor(0x6F, 0x6A, 0x61)
CREAM = RGBColor(0xF4, 0xEE, 0xE0)
CREAM_D = RGBColor(0xEB, 0xE3, 0xD2)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0xD4, 0x5A)
SUN = RGBColor(0xFF, 0xE2, 0x4B)
BLUE = RGBColor(0x2B, 0x59, 0xE0)

SITE = "https://legoschool.github.io/wise-ai/"

ERROR_SHEETS = 5      # 오류 제보 빈 양식 장수
REVIEW_SHEETS = 6     # 수업 후기 빈 양식 장수


# ---------------------------------------------------------------- 그리기 도구

def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CREAM
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def box(slide, x, y, w, h, fill=PAPER, line=INK, width=Pt(2.5), radius=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = width
    sh.shadow.inherit = False
    if radius:
        try:
            sh.adjustments[0] = 0.08
        except Exception:
            pass
    return sh


def text(slide, x, y, w, h, lines, size=18, bold=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.3):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = color
    return tb


def label(slide, x, y, w, txt, size=15):
    """칸 이름표. 초록 띠에 검은 글씨."""
    h = Inches(0.34)
    b = box(slide, x, y, w, h, fill=GREEN, line=INK, width=Pt(2))
    tf = b.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = INK
    return b


def write_here(slide, x, y, w, h, hint):
    """적는 칸. 흰 상자에 옅은 안내 문구를 넣어 둔다."""
    box(slide, x, y, w, h, fill=PAPER, line=INK, width=Pt(2))
    text(slide, x + Inches(0.14), y + Inches(0.1), w - Inches(0.28), h - Inches(0.2),
         hint, size=14, color=MUTED)


def photo_here(slide, x, y, w, h, hint):
    """사진 자리. 점선 느낌의 굵은 테두리와 안내 문구."""
    box(slide, x, y, w, h, fill=CREAM_D, line=INK, width=Pt(2.5))
    text(slide, x, y + h / 2 - Inches(0.42), w, Inches(0.9), hint,
         size=15, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def head(slide, kicker, title, tint=GREEN):
    """쪽 머리. 왼쪽에 색 띠를 세우고 제목을 적는다."""
    box(slide, Inches(0.55), Inches(0.42), Inches(0.16), Inches(0.95), fill=tint, line=None)
    text(slide, Inches(0.85), Inches(0.4), Inches(9.5), Inches(0.28), kicker,
         size=14, bold=True, color=MUTED)
    text(slide, Inches(0.85), Inches(0.72), Inches(11.5), Inches(0.66), title,
         size=28, bold=True)


def foot(slide, txt):
    text(slide, Inches(0.55), Inches(6.92), Inches(12.2), Inches(0.4), txt,
         size=12, color=MUTED)


# ---------------------------------------------------------------- 쪽 만들기

def cover(prs, prog):
    s = blank(prs)
    box(s, 0, 0, W, Inches(4.5), fill=BLUE, line=None)
    text(s, Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.4),
         prog["team"], size=16, bold=True, color=PAPER)
    text(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(2.0),
         ["써 보고 알려 주세요", "오류 제보 · 수업 후기 양식"],
         size=44, bold=True, color=PAPER, spacing=1.15)
    text(s, Inches(0.9), Inches(3.35), Inches(11.5), Inches(0.5),
         "%s · 초등 5·6학년 12차시" % prog["name"], size=17, color=PAPER)
    box(s, Inches(0.9), Inches(4.95), Inches(11.5), Inches(1.35), fill=PAPER, line=INK)
    text(s, Inches(1.2), Inches(5.12), Inches(11.0), Inches(1.0),
         ["이 파일은 채워 넣는 양식입니다. 캔바에서 열어 각자 한 벌씩 복제해 쓰세요.",
          "수업 주소 : %s" % SITE],
         size=16, spacing=1.4)
    foot(s, "채워서 연구회 대표에게 보내 주시면 다음 판에 반영합니다.")
    return s


def how_to(prs):
    s = blank(prs)
    head(s, "먼저 읽어 주세요", "이 파일을 쓰는 법")
    steps = [
        ("1", "복제합니다", "캔바에서 이 파일을 연 뒤 내 것으로 복제합니다.\n원본은 그대로 두세요."),
        ("2", "필요한 장을 골라 씁니다",
         "오류를 만나면 오류 제보 장을,\n수업을 마치면 수업 후기 장을 씁니다."),
        ("3", "모자라면 장을 복제합니다",
         "빈 양식이 모자라면 그 장을 통째로 복제해서\n계속 쓰시면 됩니다."),
    ]
    x = Inches(0.75)
    for no, title, body in steps:
        box(s, x, Inches(1.85), Inches(3.85), Inches(2.9), fill=PAPER, line=INK)
        b = box(s, x + Inches(0.28), Inches(2.1), Inches(0.62), Inches(0.62),
                fill=SUN, line=INK, width=Pt(2))
        tf = b.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = no
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.name = FONT
        r.font.color.rgb = INK
        text(s, x + Inches(0.28), Inches(2.95), Inches(3.3), Inches(0.5), title, size=20, bold=True)
        text(s, x + Inches(0.28), Inches(3.45), Inches(3.3), Inches(1.1),
             body.split("\n"), size=15, color=MUTED)
        x += Inches(4.05)
    box(s, Inches(0.75), Inches(5.05), Inches(11.85), Inches(1.55), fill=CREAM_D, line=INK, width=Pt(2))
    text(s, Inches(1.05), Inches(5.25), Inches(11.3), Inches(1.2),
         ["작은 것도 알려 주세요.",
          "글자가 어색하다, 단추를 못 찾겠다, 아이들이 이 화면에서 자꾸 멈춘다 같은 것이 "
          "가장 도움이 됩니다. 큰 오류만 적으실 필요가 없습니다."],
         size=16, spacing=1.4)
    foot(s, "오류 제보 %d장 · 수업 후기 %d장이 들어 있습니다. 모자라면 복제해서 쓰세요."
         % (ERROR_SHEETS, REVIEW_SHEETS))
    return s


def error_sheet(prs, no=None, sample=False):
    s = blank(prs)
    kicker = "오류 제보" + ("" if no is None else " %d" % no)
    head(s, kicker, "무엇이 어떻게 잘못되었나요" if not sample else "이렇게 적어 주세요 (예시)",
         tint=SUN if sample else GREEN)

    top = Inches(1.75)
    # 윗줄 : 언제 어디서
    cols = [("날짜", Inches(0.75), Inches(1.9), "예) 9월 3일"),
            ("학급", Inches(2.75), Inches(2.1), "예) 6학년 1반"),
            ("차시", Inches(5.0), Inches(1.7), "예) 3차시"),
            ("어디에서", Inches(6.85), Inches(5.75), "예) 3차시 웹앱 · 편향 관찰 갤러리 화면")]
    for name, x, w, hint in cols:
        label(s, x, top, w, name)
        write_here(s, x, top + Inches(0.38), w, Inches(0.62),
                   hint if not sample else {
                       "날짜": "9월 3일",
                       "학급": "6학년 1반",
                       "차시": "3차시",
                       "어디에서": "3차시 웹앱, 사진을 고르는 화면",
                   }[name])

    # 화면 캡처 자리
    label(s, Inches(0.75), Inches(2.95), Inches(5.9), "화면 캡처를 여기에 붙여 주세요")
    photo_here(s, Inches(0.75), Inches(3.33), Inches(5.9), Inches(3.1),
               ["이 상자 위에 캡처 그림을 끌어다 놓으세요",
                "무엇이 잘못되었는지 화살표나 동그라미로 표시해 주시면 더 좋습니다"])

    # 오른쪽 : 무슨 일이 있었나
    rx, rw = Inches(6.85), Inches(5.75)
    label(s, rx, Inches(2.95), rw, "무엇을 하려던 중이었나요")
    write_here(s, rx, Inches(3.33), rw, Inches(0.85),
               "" if not sample else "학생이 사진을 고르고 다음으로 넘어가려던 참이었습니다.")
    label(s, rx, Inches(4.3), rw, "무슨 일이 일어났나요")
    write_here(s, rx, Inches(4.68), rw, Inches(0.85),
               "" if not sample else "다음 단추를 눌러도 화면이 넘어가지 않았습니다.")
    label(s, rx, Inches(5.65), rw, "다시 해도 똑같았나요 · 무엇으로 열었나요")
    write_here(s, rx, Inches(6.03), rw, Inches(0.78),
               "" if not sample else "네, 세 번 다 그랬습니다. 학교 태블릿 크롬으로 열었습니다.")
    foot(s, "캡처는 태블릿에서 전원+음량아래 를 함께 누르면 찍힙니다. "
            "컴퓨터는 윈도우 키+Shift+S 입니다.")
    return s


def review_sheet(prs, no=None, sample=False):
    s = blank(prs)
    kicker = "수업 후기" + ("" if no is None else " %d" % no)
    head(s, kicker, "오늘 수업은 어땠나요" if not sample else "이렇게 적어 주세요 (예시)",
         tint=SUN if sample else BLUE)

    top = Inches(1.75)
    for name, x, w, hint in [("날짜", Inches(0.75), Inches(1.9), "예) 9월 3일"),
                             ("학급", Inches(2.75), Inches(2.1), "예) 6학년 1반"),
                             ("차시", Inches(5.0), Inches(1.7), "예) 3차시"),
                             ("오늘의 한 줄", Inches(6.85), Inches(5.75),
                              "예) 아이들이 편향을 눈으로 보고 놀랐습니다")]:
        label(s, x, top, w, name)
        write_here(s, x, top + Inches(0.38), w, Inches(0.62),
                   hint if not sample else {
                       "날짜": "9월 3일", "학급": "6학년 1반", "차시": "3차시",
                       "오늘의 한 줄": "데이터를 줄였더니 정확도가 무너지는 걸 보고 다들 조용해졌습니다.",
                   }[name])

    # 사진 세 자리
    label(s, Inches(0.75), Inches(2.95), Inches(11.85), "수업 사진을 붙여 주세요 (얼굴이 드러나지 않게 부탁드립니다)")
    px = Inches(0.75)
    for i in range(3):
        photo_here(s, px, Inches(3.33), Inches(3.85), Inches(2.15),
                   "사진 %d" % (i + 1))
        px += Inches(4.0)

    # 아래 세 칸
    bx, bw = Inches(0.75), Inches(3.85)
    for name, hint_sample in [("좋았던 점", "직접 라벨을 붙여 보게 한 것이 좋았습니다."),
                              ("아이들 반응", "빠른 답이 맞는 답은 아니라는 말을 아이가 먼저 했습니다."),
                              ("다음에 바꿀 점", "모둠을 네 명에서 세 명으로 줄이겠습니다.")]:
        label(s, bx, Inches(5.62), bw, name)
        write_here(s, bx, Inches(6.0), bw, Inches(0.85), "" if not sample else hint_sample)
        bx += Inches(4.0)
    foot(s, "사진은 학생 얼굴이 드러나지 않게 찍어 주세요. 손, 화면, 활동지 중심이면 좋습니다.")
    return s


def closing(prs, prog):
    s = blank(prs)
    box(s, 0, 0, W, H, fill=BLUE, line=None)
    text(s, Inches(1.0), Inches(1.6), Inches(11.3), Inches(1.2),
         ["고맙습니다"], size=48, bold=True, color=PAPER)
    text(s, Inches(1.0), Inches(3.0), Inches(11.3), Inches(2.0),
         ["선생님이 적어 주신 한 줄이 다음 판을 만듭니다.",
          "오류는 고치고, 후기는 자료에 담아 다른 선생님께 전하겠습니다.",
          "",
          "수업 주소 : %s" % SITE],
         size=19, color=PAPER, spacing=1.6)
    text(s, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.5),
         prog["copyrightLine"], size=14, color=PAPER)
    return s


# ---------------------------------------------------------------- 만들기

def build():
    data = T.load_lessons()
    prog = data["program"]
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    cover(prs, prog)
    how_to(prs)
    error_sheet(prs, sample=True)
    for i in range(1, ERROR_SHEETS + 1):
        error_sheet(prs, no=i)
    review_sheet(prs, sample=True)
    for i in range(1, REVIEW_SHEETS + 1):
        review_sheet(prs, no=i)
    closing(prs, prog)

    out_dir = os.path.join(T.ROOT, "out", "서류")
    if not os.path.isdir(out_dir):
        os.makedirs(out_dir)
    path = os.path.join(out_dir, "WISE_오류제보_수업후기_양식.pptx")
    prs.save(path)
    print("만들었다 : %s  (슬라이드 %d장)" % (path, len(prs.slides.__iter__.__self__._sldIdLst)))
    return path


def main():
    build()
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
