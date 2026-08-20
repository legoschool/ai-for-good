# -*- coding: utf-8 -*-
"""차시 수업용 PPTX 를 만든다. 발표자 노트에 교사 발문을 그대로 넣는다.

사용법 : py -3 build/make_ppt.py 3
"""
import io
import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

import tasks as T

T.setup_console()

FONT = "맑은 고딕"
W, H = Inches(13.333), Inches(7.5)

INK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xF3, 0xF4, 0xF6)


SHOT_DIR = os.path.join("out", "site", "assets", "shots")
QR_DIR = os.path.join("out", "site", "assets", "qr")
SITE_BASE = "https://legoschool.github.io/wise-ai"


def shots_of(lid):
    """make_shots.js 가 찍어 둔 화면 목록. 없으면 빈 목록."""
    path = os.path.join(T.ROOT, SHOT_DIR, "index.json")
    if not os.path.exists(path):
        return []
    try:
        with io.open(path, encoding="utf-8") as f:
            return json.load(f).get(lid, [])
    except Exception:
        return []


def shot_path(name):
    return os.path.join(T.ROOT, SHOT_DIR, name)


def qr_png(lid):
    """QR 은 SVG 로 만들어 두었다. PPT 에는 PNG 가 필요하므로 그때그때 만든다."""
    out = os.path.join(T.ROOT, QR_DIR, "%s.png" % lid)
    if os.path.exists(out):
        return out
    try:
        import segno
        segno.make("%s/webapp/%s/index.html" % (SITE_BASE, lid), error="m").save(
            out, kind="png", scale=6, border=2, dark="#111111", light="#ffffff")
        return out
    except Exception:
        return None


def hex_to_rgb(s):
    s = s.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def text(slide, x, y, w, h, runs, size=24, bold=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.25):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = color
    return box


def note(slide, lines):
    tf = slide.notes_slide.notes_text_frame
    tf.text = lines[0] if lines else ""
    for line in lines[1:]:
        tf.add_paragraph().text = line


def turn_lines(block):
    out = ["◎ " + block["heading"]]
    for t in block.get("turns", []):
        out.append("- " + t["q"])
        for a in t.get("a", []):
            out.append("  · " + a)
    return out


def build(lesson, data):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    mod = data["modules"][lesson["module"] - 1]
    accent = hex_to_rgb(mod["color"])
    M = Inches(0.9)
    CW = W - M * 2

    # 1. 표지
    s = blank(prs)
    rect(s, 0, 0, W, Inches(0.28), accent)
    text(s, M, Inches(2.0), CW, Inches(0.6),
         "%d차시  ·  모듈%d %s" % (lesson["no"], mod["no"], mod["name"]),
         size=22, color=accent, bold=True)
    text(s, M, Inches(2.7), CW, Inches(1.6), lesson["title"], size=48, bold=True)
    text(s, M, Inches(4.5), CW, Inches(1.0), lesson["problem"], size=22, color=MUTED)
    text(s, M, H - Inches(1.0), CW, Inches(0.5),
         data["program"]["copyrightLine"], size=12, color=MUTED)
    note(s, ["오늘 배울 것을 한 문장으로 소개한다.", lesson["problem"]])

    # 2. 학습 문제
    s = blank(prs)
    rect(s, 0, Inches(2.4), W, Inches(2.7), BAND)
    text(s, M, Inches(1.5), CW, Inches(0.6), "학습 문제", size=22, color=accent, bold=True)
    text(s, M, Inches(2.9), CW, Inches(1.8), lesson["problem"], size=34, bold=True)
    note(s, ["학습 문제를 함께 소리 내어 읽는다."])

    # 3. 동기 유발
    intro = lesson["plan"]["intro"]
    for b in intro["blocks"]:
        if b["heading"] == "학습 문제 확인":
            continue
        s = blank(prs)
        rect(s, 0, 0, Inches(0.22), H, accent)
        text(s, M, Inches(0.8), CW, Inches(0.6), "도입 · " + b["heading"],
             size=20, color=accent, bold=True)
        qs = [t["q"] for t in b.get("turns", [])]
        text(s, M, Inches(1.9), CW, Inches(4.2), qs, size=28, spacing=1.6)
        note(s, turn_lines(b) + ["", "준비물"] + intro["materials"])

    # 4. 전개 활동
    for b in lesson["plan"]["develop"]["blocks"]:
        s = blank(prs)
        rect(s, 0, 0, Inches(0.22), H, accent)
        text(s, M, Inches(0.8), CW, Inches(0.6), b["heading"], size=26, color=accent, bold=True)
        qs = [t["q"] for t in b.get("turns", [])]
        text(s, M, Inches(2.0), CW, Inches(4.2), qs, size=26, spacing=1.7)
        note(s, turn_lines(b) + ["", "준비물·유의점"] + lesson["plan"]["develop"]["materials"])

    # 5. 웹앱 안내
    w = lesson["webapp"]
    s = blank(prs)
    rect(s, 0, 0, Inches(0.22), H, accent)
    text(s, M, Inches(0.8), CW, Inches(0.6), "오늘 쓰는 웹앱", size=20, color=accent, bold=True)
    text(s, M, Inches(1.6), CW, Inches(0.9), w["name"], size=38, bold=True)
    text(s, M, Inches(2.7), CW, Inches(1.0), w["purpose"], size=20, color=MUTED)
    text(s, M, Inches(4.0), CW, Inches(2.2),
         ["· " + x for x in w["screens"]], size=20, spacing=1.5)
    note(s, ["방 코드를 만들어 화면에 크게 띄운다.",
             "혼자 체험 경로도 함께 안내한다.",
             "교사 화면 : " + w.get("teacherView", "")])

    # 5-1. 웹앱 단계별 안내. steps 가 있는 차시만 만든다.
    for st in w.get("steps", []):
        s = blank(prs)
        rect(s, 0, 0, Inches(0.22), H, accent)
        rect(s, W - Inches(2.3), Inches(0.7), Inches(1.4), Inches(0.55), accent)
        text(s, W - Inches(2.3), Inches(0.72), Inches(1.4), Inches(0.5),
             "%d분" % st["minutes"], size=20, bold=True, color=PAPER, align=PP_ALIGN.CENTER)
        text(s, M, Inches(0.8), CW - Inches(2.4), Inches(0.5),
             "%d단계" % st["no"], size=18, color=accent, bold=True)
        text(s, M, Inches(1.5), CW, Inches(0.9), st["title"], size=36, bold=True)
        text(s, M, Inches(2.7), CW, Inches(1.6), st["what"], size=20, color=MUTED, spacing=1.5)
        text(s, M, Inches(4.5), CW, Inches(1.2), st["ask"], size=26, bold=True)
        note(s, ["%d단계 · %d분" % (st["no"], st["minutes"]),
                 st["what"], "", "발문 : " + st["ask"], "예상 답변 : " + st["expect"]])

    # 6. AI 미사용 대안
    s = blank(prs)
    rect(s, 0, Inches(2.3), W, Inches(2.6), BAND)
    text(s, M, Inches(1.4), CW, Inches(0.6), "기기가 없어도 함께합니다",
         size=20, color=accent, bold=True)
    text(s, M, Inches(2.8), CW, Inches(1.8), lesson["alternative"], size=24)
    note(s, ["1인 1기기가 아니어도 모두 참여하게 한다."] + lesson["cautions"])

    # 7. 정리
    close = lesson["plan"]["close"]
    for b in close["blocks"]:
        s = blank(prs)
        rect(s, 0, 0, Inches(0.22), H, accent)
        text(s, M, Inches(0.8), CW, Inches(0.6), "정리 · " + b["heading"],
             size=20, color=accent, bold=True)
        qs = [t["q"] for t in b.get("turns", [])]
        text(s, M, Inches(2.0), CW, Inches(3.6), qs, size=28, spacing=1.6)
        note(s, turn_lines(b) + ["", "준비물"] + close["materials"])

    # 7-2. 웹앱 사용법. 캡처가 있을 때만 넣는다.
    lid = lesson["id"]
    shots = shots_of(lid)
    gate = None
    steps = []
    for one in shots:
        if one.get("name") == "입장":
            gate = one
        elif one.get("name") not in ("이야기",):
            steps.append(one)

    if gate or steps:
        w = lesson["webapp"]

        s = blank(prs)
        rect(s, 0, 0, Inches(0.22), H, accent)
        text(s, M, Inches(0.8), CW, Inches(0.6), "웹앱으로 합니다",
             size=20, color=accent, bold=True)
        text(s, M, Inches(1.6), Inches(7.4), Inches(1.2), w["name"], size=40, bold=True)
        text(s, M, Inches(2.9), Inches(7.4), Inches(2.4),
             [w["purpose"],
              "수업 전에 선생님 화면 → 새 방 만들기 → 방 번호를 칠판에 적는다.",
              "학생은 QR 을 찍거나 주소를 열고 방 번호와 닉네임을 넣는다.",
              "%s/webapp/%s/" % (SITE_BASE, lid)],
             size=18, spacing=1.5)
        qr = qr_png(lid)
        if qr:
            s.shapes.add_picture(qr, Inches(9.2), Inches(2.0), height=Inches(3.0))
        note(s, ["수업 시작 전에 방을 만들어 둔다.",
                 "학생에게는 방 번호와 닉네임만 알려 준다. 이름을 묻지 않는다."])

        if gate:
            s = blank(prs)
            rect(s, 0, 0, Inches(0.22), H, accent)
            text(s, M, Inches(0.7), CW, Inches(0.6), "학생 입장 화면",
                 size=20, color=accent, bold=True)
            s.shapes.add_picture(shot_path(gate["file"]), M, Inches(1.4), height=Inches(5.4))
            text(s, Inches(7.6), Inches(1.6), Inches(5.2), Inches(3.6),
                 ["방 번호 여섯 자리와 닉네임만 넣는다.",
                  "나만 아는 숫자 네 자리는 사전과 사후 설문을 잇는 데만 쓴다.",
                  "기기가 없으면 둘러보기로 교사가 시연한다."],
                 size=20, spacing=1.6)
            note(s, ["학생 화면을 띄워 놓고 함께 입장한다."])

        # 규격은 10~24장이다. 차시마다 활동 수가 달라 앞부분 장수가 다르므로,
        # 남는 자리만큼만 화면을 넣는다. 1·2차시처럼 활동이 많은 차시는 화면이 줄어든다.
        used = len(prs.slides.__iter__.__self__._sldIdLst)
        room = 24 - used - 2            # 뒤에 교사 화면과 오늘의 배움 두 장이 더 온다
        for i, one in enumerate(steps[:max(2, min(6, room))]):
            screen = w["screens"][i] if i < len(w["screens"]) else one["name"]
            s = blank(prs)
            rect(s, 0, 0, Inches(0.22), H, accent)
            text(s, M, Inches(0.7), CW, Inches(0.6),
                 "웹앱 %d단계 · %s" % (i + 1, one["name"]), size=20, color=accent, bold=True)
            s.shapes.add_picture(shot_path(one["file"]), M, Inches(1.4), height=Inches(5.4))
            text(s, Inches(7.6), Inches(1.8), Inches(5.2), Inches(3.2),
                 [screen, "학생이 여기서 남긴 것이 활동지와 이어진다."],
                 size=22, spacing=1.6)
            note(s, ["이 화면에서 학생이 무엇을 하는지 먼저 말해 준다.",
                     "판단이 갈리는 자리에서 까닭을 묻는다."])

        s = blank(prs)
        rect(s, 0, 0, Inches(0.22), H, accent)
        text(s, M, Inches(1.0), CW, Inches(0.6), "교사 화면", size=20, color=accent, bold=True)
        text(s, M, Inches(2.0), CW, Inches(3.2),
             [w.get("teacherView", "학급 집계를 본다."),
              "결과 크게 띄우기를 누르면 학급 화면으로 함께 본다.",
              "CSV 로 내려받고, 수업이 끝나면 방을 잠근다."],
             size=24, spacing=1.6)
        note(s, ["제출이 모이면 갈린 항목부터 함께 본다."])

    # 8. 오늘의 배움
    s = blank(prs)
    rect(s, 0, 0, Inches(0.22), H, accent)
    text(s, M, Inches(1.0), CW, Inches(0.6), "오늘의 배움", size=20, color=accent, bold=True)
    know = [f["knowledge"] for f in lesson["humanSkills"]["focus"]]
    text(s, M, Inches(2.0), CW, Inches(3.2), ["· " + k for k in know], size=24, spacing=1.6)
    text(s, M, H - Inches(1.0), CW, Inches(0.5),
         data["program"]["copyrightLine"], size=12, color=MUTED)
    note(s, ["오늘 배운 것을 학생이 한 문장으로 말하게 한다."])

    return prs


def make(no):
    data = T.load_lessons()
    lesson = [l for l in data["lessons"] if l["no"] == no][0]
    prs = build(lesson, data)
    out_dir = os.path.join(T.ROOT, "out", "ppt")
    if not os.path.isdir(out_dir):
        os.makedirs(out_dir)
    out = os.path.join(out_dir, "WISE_%s_수업.pptx" % lesson["id"])
    prs.save(out)
    print("만들었다 : %s  (슬라이드 %d장)" % (out, len(prs.slides.__iter__.__self__._sldIdLst)))
    return out


def main():
    if len(sys.argv) < 2:
        print("사용법 : py -3 build/make_ppt.py <차시번호|all>")
        return 1
    if sys.argv[1] == "all":
        for n in range(1, 13):
            make(n)
        return 0
    make(int(sys.argv[1]))
    return 0


if __name__ == "__main__":
    sys.exit(main())
