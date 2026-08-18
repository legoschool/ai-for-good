# -*- coding: utf-8 -*-
"""PPTX 산출물 검증."""
import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches

import tasks as T

T.setup_console()

ERRORS = []
FONT = "맑은 고딕"
MIN_SLIDES, MAX_SLIDES = 10, 16


def err(m):
    ERRORS.append(m)


def slide_text(s):
    out = []
    for sh in s.shapes:
        if sh.has_text_frame:
            out.append(sh.text_frame.text)
    return "\n".join(out)


def notes_text(s):
    if not s.has_notes_slide:
        return ""
    return s.notes_slide.notes_text_frame.text


def main():
    if len(sys.argv) < 2:
        print("사용법 : py -3 build/validate_pptx.py <경로.pptx>")
        return 1

    path = sys.argv[1]
    if not os.path.isabs(path):
        path = os.path.join(T.ROOT, path)
    if not os.path.exists(path):
        print("실패  파일이 없다 : %s" % path)
        return 1

    try:
        prs = Presentation(path)
    except Exception as e:
        print("실패  PPTX 를 열지 못했다 : %s" % e)
        return 1

    slides = list(prs.slides)

    # 화면 비율
    if abs(prs.slide_width - Inches(13.333)) > Inches(0.05):
        err("16:9 (13.333in) 가 아니다")

    if not (MIN_SLIDES <= len(slides) <= MAX_SLIDES):
        err("슬라이드가 %d장이다. %d~%d장이어야 한다" % (len(slides), MIN_SLIDES, MAX_SLIDES))

    all_text = "\n".join(slide_text(s) for s in slides)
    all_notes = "\n".join(notes_text(s) for s in slides)

    # 폰트 지정 (교사 PC에서 깨지지 않게)
    unnamed = 0
    for s in slides:
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.name != FONT:
                        unnamed += 1
    if unnamed:
        err("폰트가 %s 로 지정되지 않은 텍스트가 %d곳 있다" % (FONT, unnamed))

    # 발표자 노트
    empty_notes = [i + 1 for i, s in enumerate(slides) if not notes_text(s).strip()]
    if empty_notes:
        err("발표자 노트가 빈 슬라이드가 있다 : %s" % empty_notes)

    # 내용 일치
    m = re.search(r"WISE_L(\d\d)_", os.path.basename(path))
    if m:
        no = int(m.group(1))
        data = T.load_lessons()
        lesson = [l for l in data["lessons"] if l["no"] == no][0]

        if lesson["problem"] not in all_text:
            err("학습 문제 슬라이드가 없다")
        if lesson["title"] not in all_text and lesson["shortTitle"] not in all_text:
            err("학습 주제가 없다")
        if lesson["webapp"]["name"] not in all_text:
            err("웹앱 안내 슬라이드가 없다 : %s" % lesson["webapp"]["name"])
        if lesson["alternative"] not in all_text:
            err("AI 미사용 대안 활동 슬라이드가 없다")
        if data["program"]["copyrightLine"] not in all_text:
            err("저작권 표기가 없다")

        # 교사 발문이 노트에 들어갔는가
        missing = []
        for stage in ("intro", "develop", "close"):
            for b in lesson["plan"][stage]["blocks"]:
                for t in b.get("turns", []):
                    if t["q"] not in all_notes and t["q"] not in all_text:
                        missing.append(t["q"][:24])
        if missing:
            err("발표자 노트에 빠진 교사 발문 %d개 : %s" % (len(missing), missing[:3]))

        # 다른 차시 오염
        for other in data["lessons"]:
            if other["no"] != no and other["problem"] in all_text:
                err("다른 차시(%d) 내용이 섞여 있다" % other["no"])

    if "—" in all_text or "—" in all_notes:
        err("em dash(—) 를 쓰지 않는다")

    if ERRORS:
        for e in ERRORS:
            print("실패  %s" % e)
        print("")
        print("NG  오류 %d건" % len(ERRORS))
        return 1

    print("OK  %s" % os.path.basename(path))
    print("    슬라이드 %d장, 16:9, 발표자 노트 전부 채움, 폰트 %s" % (len(slides), FONT))
    print("    학습문제·웹앱·대안활동·저작권·교사발문 일치 확인")
    return 0


if __name__ == "__main__":
    sys.exit(main())
